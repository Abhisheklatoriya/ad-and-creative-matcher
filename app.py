import streamlit as st
from pptx import Presentation
import re
import zipfile
import io
import os

# 1. Page Configuration
st.set_page_config(page_title="Ad Matcher Session Sync", layout="wide")

# Ensure the library is imported correctly for PPTX
try:
    from pptx import Presentation
except ImportError:
    st.error("Library 'python-pptx' missing. Please check requirements.txt")

@st.cache_data
def extract_data_from_pptx(file_bytes):
    """Parses PPTX bytes to find Ad Codes and associated text blocks."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        full_slides_data = []
        all_codes = []

        for slide in prs.slides:
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
            combined_text = "\n".join(slide_text)
            # Standard MediaRadar 8-digit Ad Code regex
            codes = re.findall(r'\b\d{8}\b', combined_text) 
            
            if codes:
                all_codes.extend(codes)
                full_slides_data.append({"codes": codes, "text": combined_text})
                
        return sorted(list(set(all_codes))), full_slides_data
    except Exception as e:
        st.error(f"Error reading PowerPoint: {e}")
        return [], []

@st.cache_resource
def load_all_data(uploaded_files):
    """Handles Session Zips, individual files, and the PPTX report."""
    assets = []
    pptx_data = None
    pptx_name = None

    for uploaded_file in uploaded_files:
        fname = uploaded_file.name.lower()
        if fname.endswith('.pptx'):
            pptx_data = uploaded_file.getvalue()
            pptx_name = uploaded_file.name
        elif fname.endswith('.zip'):
            try:
                with zipfile.ZipFile(uploaded_file) as z:
                    for file_info in z.infolist():
                        if file_info.is_dir() or "__MACOSX" in file_info.filename:
                            continue
                        with z.open(file_info) as f:
                            content = f.read()
                            if file_info.filename.lower().endswith('.pptx'):
                                pptx_data = content
                                pptx_name = file_info.filename
                            else:
                                assets.append({
                                    "name": os.path.basename(file_info.filename),
                                    "data": content,
                                    "ext": file_info.filename.split('.')[-1].lower()
                                })
            except Exception as e:
                st.error(f"Error unzipping {uploaded_file.name}: {e}")
        else:
            assets.append({
                "name": uploaded_file.name,
                "data": uploaded_file.getvalue(),
                "ext": fname.split('.')[-1].lower()
            })
    return pptx_data, pptx_name, assets

# --- UI LOGIC ---
st.title("📦 Ad Matcher Session Sync")

with st.sidebar:
    st.header("Upload Data")
    # Unified uploader for convenience
    files = st.file_uploader("Upload PPTX, Assets, or Session ZIP", accept_multiple_files=True)
    
    if st.button("Reset Everything"):
        st.cache_resource.clear()
        st.cache_data.clear()
        st.rerun()

if files:
    pptx_bytes, pptx_name, all_assets = load_all_data(files)

    if pptx_bytes:
        ad_codes, slides_content = extract_data_from_pptx(pptx_bytes)
        
        # Session Export Tool
        st.sidebar.header("Step 3: Save Session")
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as z:
            z.writestr(pptx_name, pptx_bytes)
            for a in all_assets:
                # Store assets in a subfolder within the zip
                z.writestr(f"assets/{a['name']}", a['data'])
        
        st.sidebar.download_button(
            label="💾 Download Session Zip",
            data=buf.getvalue(),
            file_name="Ad_Session_Export.zip",
            mime="application/zip"
        )

        st.success(f"Session Active: {len(ad_codes)} Ads Found | {len(all_assets)} Creatives Loaded")
        
        search = st.text_input("🔍 Search Ad Code", placeholder="e.g. 48735259")
        display_codes = [c for c in ad_codes if search in c] if search else ad_codes

        for code in display_codes:
            matches = [a for a in all_assets if code in a['name']]
            relevant_text = next((item['text'] for item in slides_content if code in item['codes']), "")

            with st.expander(f"📦 Ad {code} ({len(matches)} files)", expanded=True):
                col1, col2 = st.columns([1, 1.5])
                with col1:
                    st.markdown("**PPTX Details:**")
                    st.info(relevant_text if relevant_text else "No details found for this code.")
                with col2:
                    if matches:
                        for asset in matches:
                            st.caption(f"File: {asset['name']}")
                            ext = asset['ext']
                            if ext in ['mp4', 'mov', 'webm']: st.video(asset['data'])
                            elif ext in ['mp3', 'wav']: st.audio(asset['data'])
                            elif ext in ['jpg', 'jpeg', 'png', 'gif']: st.image(asset['data'])
                            st.divider()
                    else:
                        st.warning("No creative file found matching this Ad Code.")
    else:
        st.warning("No PowerPoint file detected. Please upload the .pptx report.")
else:
    st.info("Upload your MediaRadar PPTX and creative files (or a previous session zip) to start.")
